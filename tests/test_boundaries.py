from pathlib import Path
from io import BytesIO
import time

import pytest

from rag_demo.embeddings import LocalHashEmbeddingProvider
from rag_demo.llm import LocalBoundaryLLMProvider
from rag_demo.models import AccessContext, ArtifactRecord, Chunk, Document, DocumentIn, KnowledgeBase, KnowledgeBaseCreate
from rag_demo.service import RagService
from rag_demo.store import JsonStore, SqliteStore
from rag_demo.vector_store import create_chunk_store


@pytest.fixture
def service(tmp_path: Path) -> RagService:
    return RagService(
        store=JsonStore(tmp_path),
        embeddings=LocalHashEmbeddingProvider(),
        llm=LocalBoundaryLLMProvider(),
        artifact_dir=tmp_path / "artifacts",
    )


class MissingEmbeddingProvider(LocalHashEmbeddingProvider):
    async def embed(self, texts: list[str]) -> list[list[float]]:
        embeddings = await super().embed(texts)
        return embeddings[:-1]


@pytest.mark.asyncio
async def test_query_only_returns_sources_from_selected_knowledge_base(service: RagService) -> None:
    service.create_knowledge_base(
        KnowledgeBaseCreate(id="kb_a", name="A 知识库", allowed_skills=["answer_question", "write_document"])
    )
    service.create_knowledge_base(KnowledgeBaseCreate(id="kb_b", name="B 知识库"))

    await service.add_document(
        kb_id="kb_a",
        data=DocumentIn(title="A 报销制度", content="A 知识库规定：报销需要在30天内提交发票。"),
    )
    await service.add_document(
        kb_id="kb_b",
        data=DocumentIn(title="B 采购制度", content="B 知识库规定：采购审批需要部门负责人确认。"),
    )

    response = await service.query(kb_id="kb_a", question="报销多久内提交发票？", top_k=5)

    assert response.sources
    assert {source.knowledge_base_id for source in response.sources} == {"kb_a"}
    assert all("采购审批" not in source.text for source in response.sources)


@pytest.mark.asyncio
async def test_skill_must_be_allowed_by_knowledge_base(service: RagService) -> None:
    service.create_knowledge_base(KnowledgeBaseCreate(id="kb_b", name="B 知识库", allowed_skills=["answer_question"]))
    await service.add_document(kb_id="kb_b", data=DocumentIn(title="B 资料", content="B 只允许问答。"))

    with pytest.raises(PermissionError):
        await service.run_skill(kb_id="kb_b", skill_name="write_document", instruction="写一份说明", top_k=3)


@pytest.mark.asyncio
async def test_write_document_uses_current_knowledge_base_sources(service: RagService) -> None:
    service.create_knowledge_base(KnowledgeBaseCreate(id="kb_a", name="A 知识库", allowed_skills=["write_document"]))
    await service.add_document(kb_id="kb_a", data=DocumentIn(title="A 模板", content="A 文档写作要求：标题简短，列出依据。"))

    response = await service.run_skill(kb_id="kb_a", skill_name="write_document", instruction="写一个制度摘要", top_k=3)

    assert response.skill == "write_document"
    assert response.sources
    assert {source.knowledge_base_id for source in response.sources} == {"kb_a"}


@pytest.mark.asyncio
async def test_permission_tags_filter_protected_documents(service: RagService) -> None:
    service.create_knowledge_base(KnowledgeBaseCreate(id="kb_hr", name="HR 知识库"))
    await service.add_document(
        kb_id="kb_hr",
        data=DocumentIn(title="公开制度", content="公开制度：办公用品可以按需申请。"),
    )
    await service.add_document(
        kb_id="kb_hr",
        data=DocumentIn(
            title="薪酬制度",
            content="薪酬制度：年终奖按照绩效等级计算。",
            permission_tags=["salary"],
            access=AccessContext(permission_tags=["salary"]),
        ),
    )

    summaries = service.list_documents(kb_id="kb_hr", access=AccessContext(permission_tags=["salary"]))
    assert {summary.title: summary.chunk_count for summary in summaries}["薪酬制度"] > 0

    public_response = await service.query(
        kb_id="kb_hr",
        question="年终奖怎么算？",
        top_k=5,
        access=AccessContext(permission_tags=[]),
    )
    privileged_response = await service.query(
        kb_id="kb_hr",
        question="年终奖怎么算？",
        top_k=5,
        access=AccessContext(permission_tags=["salary"]),
    )

    assert all(source.title != "薪酬制度" for source in public_response.sources)
    assert any(source.title == "薪酬制度" for source in privileged_response.sources)


@pytest.mark.asyncio
async def test_document_delete_and_reindex(service: RagService) -> None:
    service.create_knowledge_base(KnowledgeBaseCreate(id="kb_docs", name="文档库"))
    document = await service.add_document(
        kb_id="kb_docs",
        data=DocumentIn(title="长文档", content="第一段内容。" * 120),
    )

    documents = service.list_documents(kb_id="kb_docs")
    assert documents[0].chunk_count > 0

    service.chunk_store.replace_document_chunks(kb_id="kb_docs", document_id=document.id, chunks=[])
    assert service.list_documents(kb_id="kb_docs")[0].chunk_count == 0

    reindexed = await service.reindex_document(kb_id="kb_docs", document_id=document.id)
    assert reindexed.chunk_count > 0

    service.delete_document(kb_id="kb_docs", document_id=document.id)
    assert service.list_documents(kb_id="kb_docs") == []
    assert service.chunk_store.list_chunks(kb_id="kb_docs") == []


@pytest.mark.asyncio
async def test_chunks_are_stored_in_local_milvus_lite(tmp_path: Path) -> None:
    service = RagService(
        store=JsonStore(tmp_path),
        embeddings=LocalHashEmbeddingProvider(),
        llm=LocalBoundaryLLMProvider(),
        artifact_dir=tmp_path / "artifacts",
    )
    service.create_knowledge_base(KnowledgeBaseCreate(id="kb_milvus", name="Milvus 向量库"))

    await service.add_document(
        kb_id="kb_milvus",
        data=DocumentIn(title="Milvus 文档", content="Milvus Lite 应该保存本地 chunk 向量。"),
    )

    chunks = service.chunk_store.list_chunks(kb_id="kb_milvus")
    assert chunks
    assert chunks[0].embedding
    assert (tmp_path / "milvus_lite.db").exists()
    assert not (tmp_path / "chunks.json").exists()


def test_milvus_lite_stores_mixed_embedding_dimensions_in_separate_collections(tmp_path: Path) -> None:
    chunk_store = create_chunk_store(uri=tmp_path / "milvus_lite.db", collection_name="test_chunks")
    chunks = [
        Chunk(
            id="chunk_dim2",
            knowledge_base_id="kb_mixed",
            document_id="doc_dim2",
            title="二维向量",
            text="二维向量内容",
            embedding=[1.0, 0.0],
        ),
        Chunk(
            id="chunk_dim3",
            knowledge_base_id="kb_mixed",
            document_id="doc_dim3",
            title="三维向量",
            text="三维向量内容",
            embedding=[1.0, 0.0, 0.0],
        ),
    ]

    chunk_store.upsert_chunks(chunks)

    assert {chunk.id for chunk in chunk_store.list_chunks(kb_id="kb_mixed")} == {"chunk_dim2", "chunk_dim3"}
    matches = chunk_store.search_chunks(
        kb_id="kb_mixed",
        tenant_id="default",
        query_embedding=[1.0, 0.0],
        limit=5,
    )
    assert [match.chunk.id for match in matches] == ["chunk_dim2"]


def test_sqlite_store_migrates_legacy_json_metadata(tmp_path: Path) -> None:
    legacy = JsonStore(tmp_path)
    legacy.create_knowledge_base(
        KnowledgeBase(
            id="kb_legacy",
            name="旧数据知识库",
            allowed_skills=["answer_question", "write_document"],
        )
    )
    legacy.add_document(
        Document(
            id="doc_legacy",
            knowledge_base_id="kb_legacy",
            title="旧文档",
            content="旧 JSON 文档内容",
        )
    )
    legacy.add_artifact(
        ArtifactRecord(
            id="artifact_legacy",
            knowledge_base_id="kb_legacy",
            filename="legacy.md",
            media_type="text/markdown",
            skill="write_markdown",
        )
    )

    store = SqliteStore(tmp_path / "boundaryrag.sqlite3", legacy_data_dir=tmp_path)

    assert store.get_knowledge_base("kb_legacy") is not None
    assert store.get_document(kb_id="kb_legacy", document_id="doc_legacy") is not None
    assert store.get_artifact(kb_id="kb_legacy", artifact_id="artifact_legacy") is not None
    assert any(event.event_type == "storage.legacy_json_imported" for event in store.list_operation_events())


@pytest.mark.asyncio
async def test_sqlite_store_persists_operation_events(tmp_path: Path) -> None:
    store = SqliteStore(tmp_path / "boundaryrag.sqlite3")
    service = RagService(
        store=store,
        embeddings=LocalHashEmbeddingProvider(),
        llm=LocalBoundaryLLMProvider(),
        artifact_dir=tmp_path / "artifacts",
    )

    service.create_knowledge_base(KnowledgeBaseCreate(id="kb_ops", name="操作记录库"))
    await service.add_document(kb_id="kb_ops", data=DocumentIn(title="操作文档", content="操作记录应该落库。"))

    events = store.list_operation_events(limit=20)
    event_types = {event.event_type for event in events}
    assert {"knowledge_base.created", "document.created", "document.indexed"}.issubset(event_types)


@pytest.mark.asyncio
async def test_artifact_skills_generate_files(service: RagService, tmp_path: Path) -> None:
    service.create_knowledge_base(
        KnowledgeBaseCreate(
            id="kb_artifacts",
            name="生成库",
            allowed_skills=["write_markdown", "write_word", "write_ppt"],
        )
    )
    await service.add_document(kb_id="kb_artifacts", data=DocumentIn(title="素材", content="素材要求：输出要清晰。"))

    generated = []
    for skill_name, extension in [
        ("write_markdown", ".md"),
        ("write_word", ".docx"),
        ("write_ppt", ".pptx"),
    ]:
        response = await service.run_skill(
            kb_id="kb_artifacts",
            skill_name=skill_name,
            instruction="生成一份说明",
            top_k=3,
        )

        assert response.artifact is not None
        generated.append(response.artifact)
        assert response.artifact.filename.endswith(extension)
        assert " " not in response.artifact.download_url
        assert (tmp_path / "artifacts" / response.artifact.filename).exists()

    summaries = service.list_artifacts(kb_id="kb_artifacts")
    assert {summary.instruction for summary in summaries} == {"生成一份说明"}

    preview = service.preview_artifact(kb_id="kb_artifacts", artifact_id=generated[0].id)
    assert preview.instruction == "生成一份说明"
    assert "生成一份说明" in preview.content

    first_path = tmp_path / "artifacts" / generated[0].filename
    service.delete_artifact(kb_id="kb_artifacts", artifact_id=generated[0].id)
    assert not first_path.exists()
    assert all(summary.id != generated[0].id for summary in service.list_artifacts(kb_id="kb_artifacts"))


@pytest.mark.asyncio
async def test_artifact_download_requires_source_permissions(service: RagService) -> None:
    service.create_knowledge_base(
        KnowledgeBaseCreate(id="kb_secure_artifacts", name="安全生成库", allowed_skills=["write_markdown"])
    )
    await service.add_document(
        kb_id="kb_secure_artifacts",
        data=DocumentIn(
            title="薪酬素材",
            content="薪酬素材：年终奖说明。",
            permission_tags=["salary"],
            access=AccessContext(permission_tags=["salary"]),
        ),
    )

    response = await service.run_skill(
        kb_id="kb_secure_artifacts",
        skill_name="write_markdown",
        instruction="生成薪酬说明",
        top_k=3,
        access=AccessContext(permission_tags=["salary"]),
    )

    assert response.artifact is not None
    with pytest.raises(PermissionError):
        service.get_artifact_file(
            kb_id="kb_secure_artifacts",
            artifact_id=response.artifact.id,
            access=AccessContext(permission_tags=[]),
        )

    artifact, path = service.get_artifact_file(
        kb_id="kb_secure_artifacts",
        artifact_id=response.artifact.id,
        access=AccessContext(permission_tags=["salary"]),
    )
    assert artifact.permission_tags == ["salary"]
    assert path.exists()


@pytest.mark.asyncio
async def test_uploaded_md_docx_and_pptx_are_ingested(service: RagService) -> None:
    service.create_knowledge_base(KnowledgeBaseCreate(id="kb_uploads", name="上传库"))

    uploads = [
        ("upload.md", "text/markdown", "# Markdown\n\nMD 上传内容：Alpha。".encode("utf-8")),
        (
            "upload.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            _docx_bytes("Word 上传内容：Beta。"),
        ),
        (
            "upload.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            _pptx_bytes("PPT 上传内容：Gamma。"),
        ),
    ]

    for filename, content_type, data in uploads:
        await service.add_uploaded_document(
            kb_id="kb_uploads",
            filename=filename,
            content_type=content_type,
            data=data,
            permission_tags=[],
        )

    documents = service.list_documents(kb_id="kb_uploads")
    assert {document.title for document in documents} == {"upload"}
    assert sum(document.chunk_count for document in documents) == 3

    response = await service.query(kb_id="kb_uploads", question="Gamma 是什么内容？", top_k=5)
    assert any("Gamma" in source.text for source in response.sources)


@pytest.mark.asyncio
async def test_embedding_count_mismatch_fails_without_partial_index(tmp_path: Path) -> None:
    service = RagService(
        store=JsonStore(tmp_path),
        embeddings=MissingEmbeddingProvider(),
        llm=LocalBoundaryLLMProvider(),
        artifact_dir=tmp_path / "artifacts",
    )
    service.create_knowledge_base(KnowledgeBaseCreate(id="kb_bad_embeddings", name="坏向量库"))

    with pytest.raises(RuntimeError, match="embedding provider returned"):
        await service.add_document(
            kb_id="kb_bad_embeddings",
            data=DocumentIn(title="不完整索引", content="这段内容不应被部分入库。"),
        )

    documents = service.list_documents(kb_id="kb_bad_embeddings")
    assert len(documents) == 1
    assert documents[0].status == "failed"
    assert "embedding provider returned" in documents[0].error
    assert documents[0].chunk_count == 0
    assert service.chunk_store.list_chunks(kb_id="kb_bad_embeddings") == []


@pytest.mark.asyncio
async def test_document_content_limit_blocks_oversized_documents(tmp_path: Path) -> None:
    service = RagService(
        store=JsonStore(tmp_path),
        embeddings=LocalHashEmbeddingProvider(),
        llm=LocalBoundaryLLMProvider(),
        artifact_dir=tmp_path / "artifacts",
        max_document_chars=10,
    )
    service.create_knowledge_base(KnowledgeBaseCreate(id="kb_limits", name="限制库"))

    with pytest.raises(ValueError, match="too large"):
        await service.add_document(
            kb_id="kb_limits",
            data=DocumentIn(title="超长文档", content="超过十个字符的文档内容"),
        )

    assert service.list_documents(kb_id="kb_limits") == []


@pytest.mark.asyncio
async def test_upload_security_rejects_invalid_office_archive(service: RagService) -> None:
    service.create_knowledge_base(KnowledgeBaseCreate(id="kb_bad_upload", name="坏上传库"))

    with pytest.raises(ValueError, match="valid zip"):
        await service.add_uploaded_document(
            kb_id="kb_bad_upload",
            filename="not-a-docx.docx",
            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            data=b"not a zip archive",
            permission_tags=[],
        )

    assert service.list_documents(kb_id="kb_bad_upload") == []


@pytest.mark.asyncio
async def test_upload_parse_timeout_is_reported(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    import rag_demo.service as service_module

    original_parse = service_module.parse_uploaded_document

    def slow_parse_uploaded_document(**kwargs: object) -> object:
        time.sleep(0.05)
        return original_parse(**kwargs)

    monkeypatch.setattr(service_module, "parse_uploaded_document", slow_parse_uploaded_document)
    service = RagService(
        store=JsonStore(tmp_path),
        embeddings=LocalHashEmbeddingProvider(),
        llm=LocalBoundaryLLMProvider(),
        artifact_dir=tmp_path / "artifacts",
        upload_parse_timeout_seconds=0.001,
    )
    service.create_knowledge_base(KnowledgeBaseCreate(id="kb_timeout", name="超时库"))

    with pytest.raises(ValueError, match="timed out"):
        await service.add_uploaded_document(
            kb_id="kb_timeout",
            filename="slow.md",
            content_type="text/markdown",
            data=b"# slow",
            permission_tags=[],
        )

    assert service.list_documents(kb_id="kb_timeout") == []


def _docx_bytes(text: str) -> bytes:
    from docx import Document as DocxDocument

    buffer = BytesIO()
    document = DocxDocument()
    document.add_paragraph(text)
    document.save(buffer)
    return buffer.getvalue()


def _pptx_bytes(text: str) -> bytes:
    from pptx import Presentation

    buffer = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "上传 PPT"
    slide.placeholders[1].text = text
    presentation.save(buffer)
    return buffer.getvalue()
