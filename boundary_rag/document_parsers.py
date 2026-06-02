"""上传文件解析和轻量安全检查。

上传文件会先转换成纯文本，然后再进入切分流程。Markdown 按文本读取；
Office 文件会先按 zip 归档检查，用于防护路径穿越、加密文件、zip bomb
以及扩展名伪装。解析器刻意只返回规范化文本和 metadata，不把上传二进制
直接存入数据库。
"""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path, PurePosixPath
from zipfile import BadZipFile, ZipFile


@dataclass(frozen=True)
class ParsedDocument:
    """上传解析器返回的规范化文档文本。"""
    title: str
    content: str
    metadata: dict[str, str]


@dataclass(frozen=True)
class UploadSecurityPolicy:
    """解析前扫描 Office zip 容器时使用的安全限制。"""
    max_archive_members: int = 512
    max_archive_uncompressed_bytes: int = 100 * 1024 * 1024
    max_archive_compression_ratio: float = 100.0


def parse_uploaded_document(
    *,
    filename: str,
    content_type: str | None,
    data: bytes,
    security_policy: UploadSecurityPolicy | None = None,
) -> ParsedDocument:
    """把受支持的上传文件解析成文本和来源 metadata。"""
    if not data:
        raise ValueError("uploaded file is empty")

    safe_filename = Path(filename).name or "uploaded-document"
    path = Path(safe_filename)
    extension = path.suffix.lower()
    title = path.stem or "uploaded-document"
    security_policy = security_policy or UploadSecurityPolicy()
    _validate_upload_security(extension=extension, data=data, policy=security_policy)

    if extension in {".md", ".markdown"}:
        content = _decode_text(data)
        parser = "markdown"
    elif extension == ".docx":
        content = _parse_docx(data)
        parser = "docx"
    elif extension == ".pptx":
        content = _parse_pptx(data)
        parser = "pptx"
    else:
        raise ValueError("unsupported file type; please upload .md, .docx, or .pptx")

    normalized = content.strip()
    if not normalized:
        raise ValueError("uploaded file contains no extractable text")

    return ParsedDocument(
        title=title,
        content=normalized,
        metadata={
            "source_filename": safe_filename,
            "content_type": content_type or "",
            "parser": parser,
        },
    )


def _validate_upload_security(*, extension: str, data: bytes, policy: UploadSecurityPolicy) -> None:
    """在提取内容前，根据扩展名分发对应的安全检查。"""
    if extension in {".md", ".markdown"}:
        if b"\x00" in data:
            raise ValueError("uploaded markdown appears to be binary data")
        return

    if extension == ".docx":
        _validate_office_zip(data=data, expected_member="word/document.xml", policy=policy)
        return

    if extension == ".pptx":
        _validate_office_zip(data=data, expected_member="ppt/presentation.xml", policy=policy)
        return

    raise ValueError("unsupported file type; please upload .md, .docx, or .pptx")


def _validate_office_zip(*, data: bytes, expected_member: str, policy: UploadSecurityPolicy) -> None:
    """在不落盘解压的情况下校验 Office 归档文件。"""
    try:
        with ZipFile(BytesIO(data)) as archive:
            infos = archive.infolist()
    except BadZipFile as exc:
        raise ValueError("uploaded Office file is not a valid zip archive") from exc

    if not infos:
        raise ValueError("uploaded Office file archive is empty")
    if len(infos) > policy.max_archive_members:
        raise ValueError("uploaded Office file contains too many archive entries")

    total_uncompressed = 0
    names: set[str] = set()
    for info in infos:
        normalized_name = info.filename.replace("\\", "/")
        parts = PurePosixPath(normalized_name).parts
        if normalized_name.startswith("/") or ".." in parts or any(":" in part for part in parts):
            raise ValueError("uploaded Office file contains unsafe archive paths")
        if info.flag_bits & 0x1:
            raise ValueError("encrypted Office files are not supported")

        total_uncompressed += info.file_size
        if total_uncompressed > policy.max_archive_uncompressed_bytes:
            raise ValueError("uploaded Office file expands beyond the safe size limit")
        if info.file_size > 0 and info.compress_size == 0:
            raise ValueError("uploaded Office file has an unsafe compression ratio")
        if info.compress_size > 0 and (info.file_size / info.compress_size) > policy.max_archive_compression_ratio:
            raise ValueError("uploaded Office file has an unsafe compression ratio")
        names.add(normalized_name)

    if expected_member not in names:
        raise ValueError("uploaded Office file does not match its file extension")


def _decode_text(data: bytes) -> str:
    """使用常见中文/UTF 编码解码 Markdown 上传内容。"""
    for encoding in ("utf-8-sig", "utf-8", "gb18030"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _parse_docx(data: bytes) -> str:
    """提取 Word 段落，并把表格扁平化为管道分隔行。"""
    from docx import Document as DocxDocument

    document = DocxDocument(BytesIO(data))
    parts: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def _parse_pptx(data: bytes) -> str:
    """提取 PPT 页面文本和表格单元格，并转换成类 Markdown 章节。"""
    from pptx import Presentation

    presentation = Presentation(BytesIO(data))
    slides: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                texts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        texts.append(" | ".join(cells))
        if texts:
            slides.append(f"## Slide {index}\n" + "\n".join(texts))

    return "\n\n".join(slides)
