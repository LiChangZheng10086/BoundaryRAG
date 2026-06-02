"""技能注册表和生成产物逻辑。

技能是知识库的能力边界。一个 KB 可以独立允许问答、写作、Markdown 导出、
Word 导出和 PPT 导出。每个技能仍然只接收当前 KB 检索出的分块，因此生成产物
和普通回答一样，都被限制在同一个 RAG 边界内。
"""

from __future__ import annotations

import re
from abc import ABC, abstractmethod
from pathlib import Path
from urllib.parse import quote
from uuid import uuid4

from boundary_rag.llm import LLMProvider
from boundary_rag.models import GeneratedArtifact, KnowledgeBase, SkillResponse, Source


class Skill(ABC):
    """所有知识库技能的基础接口。"""

    name: str

    @abstractmethod
    async def run(
        self,
        *,
        kb: KnowledgeBase,
        instruction: str,
        sources: list[Source],
        llm: LLMProvider,
    ) -> SkillResponse:
        raise NotImplementedError


class AnswerQuestionSkill(Skill):
    """默认 RAG 问答技能。"""

    name = "answer_question"

    async def run(
        self,
        *,
        kb: KnowledgeBase,
        instruction: str,
        sources: list[Source],
        llm: LLMProvider,
    ) -> SkillResponse:
        answer = await llm.answer(kb=kb, instruction=instruction, sources=sources)
        return SkillResponse(skill=self.name, answer=answer, sources=sources)


class WriteDocumentSkill(Skill):
    """生成写作结果，但不创建可下载文件。"""

    name = "write_document"

    async def run(
        self,
        *,
        kb: KnowledgeBase,
        instruction: str,
        sources: list[Source],
        llm: LLMProvider,
    ) -> SkillResponse:
        prompt = (
            f"请只根据当前知识库资料完成写作任务：{instruction}。"
            "如果资料不足，先说明缺口，再给出可依据现有资料完成的部分。"
        )
        answer = await llm.answer(kb=kb, instruction=prompt, sources=sources)
        return SkillResponse(skill=self.name, answer=answer, sources=sources)


class ArtifactSkill(Skill):
    """会在产物目录下创建文件的技能基类。"""

    media_type: str
    extension: str

    def __init__(self, artifact_dir: Path) -> None:
        self.artifact_dir = artifact_dir
        self.artifact_dir.mkdir(parents=True, exist_ok=True)

    async def run(
        self,
        *,
        kb: KnowledgeBase,
        instruction: str,
        sources: list[Source],
        llm: LLMProvider,
    ) -> SkillResponse:
        """使用 LLM 生成内容，写入文件，并返回产物 metadata。"""
        prompt = self._prompt(instruction)
        answer = await llm.answer(kb=kb, instruction=prompt, sources=sources)
        artifact_id = f"artifact_{uuid4().hex}"
        filename = self._filename(kb=kb, instruction=instruction)
        path = self.artifact_dir / filename
        self._write(path=path, title=instruction, content=answer, sources=sources)
        artifact = GeneratedArtifact(
            id=artifact_id,
            filename=filename,
            media_type=self.media_type,
            download_url=f"/knowledge-bases/{quote(kb.id)}/artifacts/{artifact_id}/download",
            instruction=instruction,
        )
        return SkillResponse(skill=self.name, answer=answer, sources=sources, artifact=artifact)

    def _prompt(self, instruction: str) -> str:
        """把用户指令包装成只允许使用当前 KB 的产物生成要求。"""
        return (
            f"请只根据当前知识库资料生成{self.extension}文档内容：{instruction}。"
            "要求结构清晰、标题明确、内容可直接交付；如果资料不足，必须说明缺口。"
        )

    def _filename(self, *, kb: KnowledgeBase, instruction: str) -> str:
        """创建安全的文件名，并保留必要的用户指令上下文。"""
        stem = re.sub(r"[^a-zA-Z0-9\u4e00-\u9fff_-]+", "-", instruction).strip("-")[:36]
        if not stem:
            stem = self.name
        return f"{kb.id}-{stem}-{uuid4().hex[:8]}.{self.extension}"

    @abstractmethod
    def _write(self, *, path: Path, title: str, content: str, sources: list[Source]) -> None:
        raise NotImplementedError


class WriteMarkdownSkill(ArtifactSkill):
    """写入 Markdown 文档产物。"""

    name = "write_markdown"
    media_type = "text/markdown"
    extension = "md"

    def _write(self, *, path: Path, title: str, content: str, sources: list[Source]) -> None:
        source_lines = "\n".join(
            f"- {source.title} (`{source.chunk_id}`), score={source.score:.3f}" for source in sources
        )
        body = f"# {title}\n\n{content}\n"
        if source_lines:
            body += f"\n## 来源\n\n{source_lines}\n"
        path.write_text(body, encoding="utf-8")


class WriteWordSkill(ArtifactSkill):
    """把生成的类 Markdown 内容写成 Word `.docx` 产物。"""

    name = "write_word"
    media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    extension = "docx"

    def _write(self, *, path: Path, title: str, content: str, sources: list[Source]) -> None:
        """把标题和项目符号转换成基础 Word 结构。"""
        from docx import Document as DocxDocument

        document = DocxDocument()
        document.add_heading(title, level=1)
        for line in content.splitlines():
            stripped = line.strip()
            if not stripped:
                continue
            if stripped.startswith("#"):
                document.add_heading(stripped.lstrip("#").strip(), level=min(stripped.count("#"), 3))
            elif stripped.startswith(("- ", "* ")):
                document.add_paragraph(stripped[2:].strip(), style="List Bullet")
            else:
                document.add_paragraph(stripped)
        if sources:
            document.add_heading("来源", level=2)
            for source in sources:
                document.add_paragraph(
                    f"{source.title} ({source.chunk_id}), score={source.score:.3f}",
                    style="List Bullet",
                )
        document.save(path)


class WritePptSkill(ArtifactSkill):
    """根据生成的章节内容写入 PowerPoint `.pptx` 产物。"""

    name = "write_ppt"
    media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    extension = "pptx"

    def _write(self, *, path: Path, title: str, content: str, sources: list[Source]) -> None:
        """创建标题页、内容页，以及可选的来源页。"""
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation()
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = title
        title_slide.placeholders[1].text = "基于当前知识库自动生成"

        sections = self._sections(content)
        for section_title, bullets in sections[:8]:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = section_title[:60]
            body = slide.placeholders[1].text_frame
            body.clear()
            for index, bullet in enumerate(bullets[:5]):
                paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
                paragraph.text = bullet[:120]
                paragraph.font.size = Pt(20)

        if sources:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "来源"
            body = slide.placeholders[1].text_frame
            body.clear()
            for index, source in enumerate(sources[:6]):
                paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
                paragraph.text = f"{source.title} · score {source.score:.3f}"
                paragraph.font.size = Pt(18)

        prs.save(path)

    def _sections(self, content: str) -> list[tuple[str, list[str]]]:
        """把生成的类 Markdown 内容解析成 PPT 章节。"""
        sections: list[tuple[str, list[str]]] = []
        current_title = "要点概览"
        current_bullets: list[str] = []

        for raw_line in content.splitlines():
            line = raw_line.strip()
            if not line:
                continue
            if line.startswith("#"):
                if current_bullets:
                    sections.append((current_title, current_bullets))
                current_title = line.lstrip("#").strip() or "要点"
                current_bullets = []
            else:
                current_bullets.append(line.lstrip("-*0123456789.、 ").strip())

        if current_bullets:
            sections.append((current_title, current_bullets))
        if not sections:
            sections.append(("要点概览", [content[:120] or "暂无内容"]))
        return sections


class SkillRegistry:
    """负责强制执行每个 KB 已配置技能的查找表。"""

    def __init__(self, artifact_dir: Path) -> None:
        skills: list[Skill] = [
            AnswerQuestionSkill(),
            WriteDocumentSkill(),
            WriteMarkdownSkill(artifact_dir),
            WriteWordSkill(artifact_dir),
            WritePptSkill(artifact_dir),
        ]
        self._skills = {skill.name: skill for skill in skills}

    def get_allowed(self, *, kb: KnowledgeBase, skill_name: str) -> Skill:
        """只有当前 KB 明确允许时才返回对应技能。"""
        if skill_name not in kb.allowed_skills:
            raise PermissionError(f"skill '{skill_name}' is not allowed for knowledge base '{kb.id}'")
        if skill_name not in self._skills:
            raise KeyError(f"unknown skill '{skill_name}'")
        return self._skills[skill_name]
